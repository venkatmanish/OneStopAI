from __future__ import annotations

import re
import uuid
from pathlib import Path
from typing import Any

from core.db import Chunk, Document, DocumentVersion, session_scope
from core.embeddings import EmbeddingModel
from core.repository import Repository
from core.schemas import IngestResult
from core.storage import LocalObjectStore, sha256_bytes


def chunk_text(text: str, size: int = 1200, overlap: int = 150) -> list[str]:
    """Split text on semantic boundaries before falling back to hard character windows."""
    cleaned = re.sub(r"\n{3,}", "\n\n", text.strip())
    if len(cleaned) <= size:
        return [cleaned] if cleaned else []

    units = _semantic_units(cleaned, size=size)
    if units:
        chunks = _pack_semantic_units(units, size=size, overlap=overlap)
        if chunks:
            return chunks

    chunks: list[str] = []
    start = 0
    while start < len(cleaned):
        end = min(len(cleaned), start + size)
        chunks.append(cleaned[start:end].strip())
        start = max(end - overlap, end) if end == len(cleaned) else end - overlap
    return [chunk for chunk in chunks if chunk]


def _semantic_units(text: str, size: int) -> list[str]:
    units: list[str] = []
    for block in re.split(r"\n\s*\n", text):
        block = block.strip()
        if not block:
            continue
        if len(block) <= size:
            units.append(block)
            continue
        units.extend(_split_long_block(block, size=size))
    return [unit for unit in units if unit]


def _split_long_block(block: str, size: int) -> list[str]:
    sentences = re.split(r"(?<=[.!?])\s+(?=[A-Z0-9])", block)
    units: list[str] = []
    current = ""
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if len(sentence) > size:
            if current:
                units.append(current.strip())
                current = ""
            units.extend(_hard_split(sentence, size=size))
            continue
        candidate = f"{current} {sentence}".strip()
        if len(candidate) <= size:
            current = candidate
        else:
            if current:
                units.append(current.strip())
            current = sentence
    if current:
        units.append(current.strip())
    return units


def _hard_split(text: str, size: int) -> list[str]:
    parts: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + size)
        if end < len(text):
            space = text.rfind(" ", start + size // 2, end)
            if space > start:
                end = space
        parts.append(text[start:end].strip())
        start = end
    return [part for part in parts if part]


def _pack_semantic_units(units: list[str], size: int, overlap: int) -> list[str]:
    chunks: list[str] = []
    current_units: list[str] = []
    current_len = 0
    for unit in units:
        separator = 2 if current_units else 0
        if current_units and current_len + separator + len(unit) > size:
            chunks.append("\n\n".join(current_units).strip())
            current_units = _overlap_units(current_units, overlap=overlap)
            current_len = sum(len(item) for item in current_units) + 2 * max(len(current_units) - 1, 0)
        current_units.append(unit)
        current_len += separator + len(unit)
    if current_units:
        chunks.append("\n\n".join(current_units).strip())
    return [chunk for chunk in chunks if chunk]


def _overlap_units(units: list[str], overlap: int) -> list[str]:
    selected: list[str] = []
    total = 0
    for unit in reversed(units):
        if selected and total + len(unit) > overlap:
            break
        selected.append(unit)
        total += len(unit)
    return list(reversed(selected))


def _compact_parent_context(text: str, limit: int = 3200) -> str:
    cleaned = re.sub(r"\s+", " ", text).strip()
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 3].rstrip() + "..."


class IngestionService:
    PARSER_VERSION = "2026-05-02-semantic-parent-rag-v1"

    def __init__(self) -> None:
        self.storage = LocalObjectStore()
        self.embedding_model = EmbeddingModel()

    def ingest_upload(self, file_name: str, content: bytes, source_type: str = "upload") -> IngestResult:
        content_hash = sha256_bytes(content)
        storage_path = self.storage.save_bytes(file_name, content, content_hash)

        with session_scope() as session:
            repo = Repository(session)
            document, version, reused = repo.upsert_document_version(
                file_name=file_name,
                source_type=source_type,
                content_hash=content_hash,
                storage_path=storage_path,
            )
            if reused:
                existing_chunks = (
                    session.query(Chunk).filter(Chunk.version_id == version.version_id).all()
                )
                if not existing_chunks or self._needs_reindex(existing_chunks):
                    extracted = self._extract(storage_path)
                    rows = self._build_chunk_rows(
                        document.document_id,
                        version.version_id,
                        file_name,
                        extracted,
                        source_type,
                    )
                    repo.replace_chunks(version.version_id, rows)
                    session.flush()
                    self._index_vectors(rows)
                    return IngestResult(
                        document_id=document.document_id,
                        version_id=version.version_id,
                        reused_existing=True,
                        chunks_indexed=len(rows),
                        summary=self._summary(file_name, len(rows), reindexed=True),
                    )
                return IngestResult(
                    document_id=document.document_id,
                    version_id=version.version_id,
                    reused_existing=True,
                    chunks_indexed=0,
                    summary="Identical content hash already indexed; existing version reused.",
                )

            extracted = self._extract(storage_path)
            rows = self._build_chunk_rows(
                document.document_id,
                version.version_id,
                file_name,
                extracted,
                source_type,
            )
            repo.replace_chunks(version.version_id, rows)
            session.flush()

            self._index_vectors(rows)
            return IngestResult(
                document_id=document.document_id,
                version_id=version.version_id,
                reused_existing=False,
                chunks_indexed=len(rows),
                summary=self._summary(file_name, len(rows)),
            )

    @staticmethod
    def _summary(file_name: str, chunk_count: int, reindexed: bool = False) -> str:
        if chunk_count == 0:
            suffix = Path(file_name).suffix.lower()
            if suffix == ".ppt":
                return (
                    f"No text chunks extracted from {file_name}. "
                    "Legacy .ppt files require LibreOffice for conversion; .pptx files are parsed directly."
                )
            if suffix == ".pptx":
                return (
                    f"No text chunks extracted from {file_name}. "
                    "Check that the slides contain selectable text, notes, tables, hyperlinks, or readable images."
                )
            return (
                f"No text chunks extracted from {file_name}. "
                "For scanned PDFs, check that the scan is readable and not heavily cropped or blurred."
            )
        prefix = "Re-indexed" if reindexed else "Indexed"
        return f"{prefix} {chunk_count} chunks from {file_name}."

    @classmethod
    def _needs_reindex(cls, chunks: list[Chunk]) -> bool:
        current_embedding_model = EmbeddingModel().settings.embedding_model
        return any(
            (chunk.extra or {}).get("parser_version") != cls.PARSER_VERSION
            or (chunk.extra or {}).get("embedding_model") != current_embedding_model
            for chunk in chunks
        )

    def reindex_active_documents(self, force: bool = True) -> dict[str, Any]:
        with session_scope() as session:
            records = (
                session.query(Document, DocumentVersion)
                .join(DocumentVersion, Document.active_version_id == DocumentVersion.version_id)
                .filter(DocumentVersion.active.is_(True))
                .all()
            )
            documents = [
                {
                    "document_id": document.document_id,
                    "version_id": version.version_id,
                    "file_name": document.file_name,
                    "source_type": document.source_type,
                    "storage_path": version.storage_path,
                }
                for document, version in records
            ]

        results: list[dict[str, Any]] = []
        total_chunks = 0
        reindexed = 0
        skipped = 0
        errors: list[str] = []
        for document in documents:
            try:
                result = self.reindex_version(
                    document_id=document["document_id"],
                    version_id=document["version_id"],
                    file_name=document["file_name"],
                    source_type=document["source_type"],
                    storage_path=Path(document["storage_path"]),
                    force=force,
                )
                results.append(result)
                if result["reindexed"]:
                    reindexed += 1
                    total_chunks += result["chunks_indexed"]
                else:
                    skipped += 1
            except Exception as exc:
                errors.append(f"{document['file_name']}: {exc}")

        return {
            "status": "ok" if not errors else "partial",
            "documents_seen": len(documents),
            "reindexed": reindexed,
            "skipped": skipped,
            "chunks_indexed": total_chunks,
            "embedding_model": self.embedding_model.settings.embedding_model,
            "parser_version": self.PARSER_VERSION,
            "results": results,
            "errors": errors,
        }

    def reindex_version(
        self,
        document_id: str,
        version_id: str,
        file_name: str,
        source_type: str,
        storage_path: Path,
        force: bool = True,
    ) -> dict[str, Any]:
        with session_scope() as session:
            existing_chunks = session.query(Chunk).filter(Chunk.version_id == version_id).all()
            if not force and existing_chunks and not self._needs_reindex(existing_chunks):
                return {
                    "document_id": document_id,
                    "version_id": version_id,
                    "file_name": file_name,
                    "reindexed": False,
                    "chunks_indexed": 0,
                    "summary": "Already indexed with the current parser and embedding model.",
                }

        extracted = self._extract(storage_path)
        rows = self._build_chunk_rows(
            document_id,
            version_id,
            file_name,
            extracted,
            source_type,
        )
        with session_scope() as session:
            Repository(session).replace_chunks(version_id, rows)
            session.flush()
        self._index_vectors(rows)
        return {
            "document_id": document_id,
            "version_id": version_id,
            "file_name": file_name,
            "reindexed": True,
            "chunks_indexed": len(rows),
            "summary": self._summary(file_name, len(rows), reindexed=True),
        }

    def _extract(self, path: Path) -> list[dict[str, Any]]:
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return PDFExtractor().extract(path)
        if suffix in {".ppt", ".pptx"}:
            return PowerPointExtractor().extract(path)
        if suffix in {".xlsx", ".xls", ".csv"}:
            return ExcelExtractor().extract(path)
        if suffix in {".png", ".jpg", ".jpeg", ".tif", ".tiff"}:
            return ChartOCRExtractor().extract(path)
        return [{"text": path.read_text(errors="ignore"), "chunk_type": "text"}]

    def _build_chunk_rows(
        self,
        document_id: str,
        version_id: str,
        file_name: str,
        extracted: list[dict[str, Any]],
        source_type: str,
    ) -> list[dict]:
        rows: list[dict] = []
        for item in extracted:
            source_text = item.get("text", "")
            chunks = chunk_text(source_text)
            parent_id = uuid.uuid4().hex
            parent_context = _compact_parent_context(source_text)
            for index, text in enumerate(chunks):
                chunk_id = uuid.uuid4().hex
                rows.append(
                    {
                        "chunk_id": chunk_id,
                        "document_id": document_id,
                        "version_id": version_id,
                        "text": text,
                        "source_type": item.get("source_type", source_type),
                        "page": item.get("page"),
                        "sheet": item.get("sheet"),
                        "chunk_type": item.get("chunk_type", "text"),
                        "extra": {
                            "file_name": file_name,
                            "table": item.get("table"),
                            "links": item.get("links"),
                            "slide": item.get("slide"),
                            "confidence": item.get("confidence"),
                            "parent_id": parent_id,
                            "parent_context": parent_context,
                            "chunk_index": index,
                            "chunk_count": len(chunks),
                            "embedding_model": self.embedding_model.settings.embedding_model,
                            "parser_version": self.PARSER_VERSION,
                        },
                    }
                )
        return rows

    def _index_vectors(self, rows: list[dict]) -> None:
        if not rows:
            return
        try:
            from core.vector_store import VectorStore

            texts = [row["text"] for row in rows]
            vectors = self.embedding_model.embed(texts)
            payloads = [
                {
                    "document_id": row["document_id"],
                    "version_id": row["version_id"],
                    "file_name": row["extra"].get("file_name"),
                    "source_type": row["source_type"],
                    "page": row["page"],
                    "sheet": row["sheet"],
                    "chunk_type": row["chunk_type"],
                    "active": True,
                }
                for row in rows
            ]
            VectorStore().upsert([row["chunk_id"] for row in rows], vectors, payloads)
        except Exception:
            # PostgreSQL chunks remain queryable by BM25 even if Qdrant is not running.
            return


class PDFExtractor:
    min_text_chars = 80

    def extract(self, path: Path) -> list[dict[str, Any]]:
        import fitz

        items: list[dict[str, Any]] = []
        doc = fitz.open(path)
        for page_index, page in enumerate(doc, start=1):
            text = self._post_process(page.get_text("text"))
            if len(text) >= self.min_text_chars:
                items.append({"text": text, "page": page_index, "chunk_type": "pdf_text"})
            else:
                ocr_text = self._extract_page_ocr(page)
                if ocr_text:
                    items.append(
                        {
                            "text": ocr_text,
                            "page": page_index,
                            "chunk_type": "ocr_text",
                            "confidence": 0.7,
                        }
                    )
                elif text:
                    items.append(
                        {
                            "text": text,
                            "page": page_index,
                            "chunk_type": "pdf_text_low_confidence",
                        }
                    )

            links = self._extract_page_links(page, page_index)
            if links:
                items.append(
                    {
                        "text": self._format_links("Hyperlinks extracted from PDF page", page_index, links),
                        "page": page_index,
                        "chunk_type": "pdf_links",
                        "links": links,
                        "confidence": 0.95,
                    }
                )
        return items

    def _extract_text_layer(self, path: Path) -> list[dict[str, Any]]:
        import fitz

        items: list[dict[str, Any]] = []
        doc = fitz.open(path)
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                items.append({"text": text, "page": page_index, "chunk_type": "pdf_text"})
        return items

    def _extract_ocr(self, path: Path) -> list[dict[str, Any]]:
        import fitz

        items: list[dict[str, Any]] = []
        doc = fitz.open(path)
        for page_index, page in enumerate(doc, start=1):
            cleaned = self._extract_page_ocr(page)
            if cleaned:
                items.append(
                    {
                        "text": cleaned,
                        "page": page_index,
                        "chunk_type": "ocr_text",
                        "confidence": 0.7,
                    }
                )
        return items

    def _extract_page_ocr(self, page) -> str:
        import cv2
        import numpy as np

        pix = page.get_pixmap(dpi=300, alpha=False)
        image = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
        if pix.n == 1:
            gray = image[:, :, 0]
        else:
            gray = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
        return self._post_process(self._best_ocr_text(gray))

    @classmethod
    def _extract_page_links(cls, page, page_index: int) -> list[dict[str, str]]:
        import fitz

        links: list[dict[str, str]] = []
        seen: set[tuple[str, str]] = set()
        for link in page.get_links():
            target = cls._link_target(link)
            if not target:
                continue
            rect = fitz.Rect(link.get("from"))
            label = cls._anchor_text_for_rect(page, rect) or "Unlabeled link"
            key = (label.lower(), target)
            if key in seen:
                continue
            seen.add(key)
            links.append(
                {
                    "label": label,
                    "target": target,
                    "kind": str(link.get("kind", "")),
                    "page": str(page_index),
                }
            )
        return links

    @classmethod
    def _anchor_text_for_rect(cls, page, rect) -> str:
        import fitz

        text = cls._post_process(page.get_textbox(rect))
        if text:
            return cls._compact_link_label(text)

        words = []
        for word in page.get_text("words"):
            word_rect = fitz.Rect(word[:4])
            if rect.intersects(word_rect):
                words.append((word[5], word_rect.y0, word_rect.x0))
        if not words:
            return ""
        words.sort(key=lambda item: (round(item[1], 1), item[2]))
        return cls._compact_link_label(" ".join(word for word, _, _ in words))

    @staticmethod
    def _link_target(link: dict[str, Any]) -> str:
        uri = link.get("uri")
        if uri:
            return str(uri).strip()
        file_name = link.get("file")
        if file_name:
            return str(file_name).strip()
        page = link.get("page")
        if page is not None and int(page) >= 0:
            return f"internal page {int(page) + 1}"
        return ""

    @staticmethod
    def _compact_link_label(label: str, limit: int = 120) -> str:
        label = re.sub(r"\s+", " ", label).strip(" :-\n\t")
        if len(label) <= limit:
            return label
        return label[: limit - 3].rstrip() + "..."

    @classmethod
    def _format_links(cls, heading: str, index: int, links: list[dict[str, str]]) -> str:
        lines = [f"{heading} {index}:"]
        for link in links:
            lines.append(f"- {link['label']} -> {link['target']}")
        return "\n".join(lines)

    @staticmethod
    def _best_ocr_text(gray) -> str:
        import cv2
        import pytesseract

        def score(text: str) -> int:
            return len(re.findall(r"[A-Za-z0-9]{2,}", text))

        blurred = cv2.medianBlur(gray, 3)
        otsu = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
        adaptive = cv2.adaptiveThreshold(
            gray,
            255,
            cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY,
            35,
            11,
        )
        candidates = [
            (adaptive, "--oem 3 --psm 6"),
            (adaptive, "--oem 3 --psm 11"),
            (otsu, "--oem 3 --psm 6"),
            (gray, "--oem 3 --psm 6"),
        ]

        best_text = ""
        best_score = 0
        for image, config in candidates:
            text = pytesseract.image_to_string(image, config=config)
            current_score = score(text)
            if current_score > best_score:
                best_text = text
                best_score = current_score
        if best_score >= 8:
            return best_text

        for rotate_code in (
            cv2.ROTATE_90_CLOCKWISE,
            cv2.ROTATE_90_COUNTERCLOCKWISE,
            cv2.ROTATE_180,
        ):
            rotated = cv2.rotate(adaptive, rotate_code)
            text = pytesseract.image_to_string(rotated, config="--oem 3 --psm 6")
            current_score = score(text)
            if current_score > best_score:
                best_text = text
                best_score = current_score
        return best_text

    @staticmethod
    def _post_process(text: str) -> str:
        text = re.sub(r"-\n(?=[a-z])", "", text)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


class PowerPointExtractor:
    def extract(self, path: Path) -> list[dict[str, Any]]:
        if path.suffix.lower() == ".ppt":
            import tempfile

            with tempfile.TemporaryDirectory() as temp_dir:
                converted = self._convert_legacy_ppt(path, Path(temp_dir))
                if converted:
                    return self._extract_pptx(converted)
            return []
        return self._extract_pptx(path)

    def _extract_pptx(self, path: Path) -> list[dict[str, Any]]:
        from pptx import Presentation

        presentation = Presentation(path)
        items: list[dict[str, Any]] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines = [f"Slide {slide_index}"]
            links: list[dict[str, str]] = []
            image_ocr: list[str] = []

            for shape in slide.shapes:
                self._collect_shape(shape, lines, links, image_ocr)

            notes = self._notes_text(slide)
            if notes:
                lines.append("Speaker notes:")
                lines.append(notes)
            if links:
                lines.append(PDFExtractor._format_links("Hyperlinks extracted from slide", slide_index, links))
            if image_ocr:
                lines.append("Image OCR:")
                lines.extend(image_ocr)

            text = PDFExtractor._post_process("\n".join(line for line in lines if line).strip())
            if text and text != f"Slide {slide_index}":
                items.append(
                    {
                        "text": text,
                        "page": slide_index,
                        "slide": slide_index,
                        "chunk_type": "ppt_slide",
                        "links": links or None,
                        "confidence": 0.9,
                    }
                )
        return items

    def _collect_shape(
        self,
        shape: Any,
        lines: list[str],
        links: list[dict[str, str]],
        image_ocr: list[str],
    ) -> None:
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                self._collect_shape(child, lines, links, image_ocr)
            return

        shape_text = self._shape_text(shape)
        if shape_text:
            lines.append(shape_text)

        self._collect_shape_hyperlinks(shape, shape_text, links)

        if getattr(shape, "has_table", False):
            table_text = self._table_text(shape)
            if table_text:
                lines.append(table_text)

        ocr_text = self._image_ocr(shape)
        if ocr_text:
            image_ocr.append(ocr_text)

    @staticmethod
    def _shape_text(shape: Any) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        text = getattr(shape, "text", "") or ""
        return PDFExtractor._post_process(text)

    @classmethod
    def _collect_shape_hyperlinks(
        cls,
        shape: Any,
        shape_text: str,
        links: list[dict[str, str]],
    ) -> None:
        seen = {(link["label"].lower(), link["target"]) for link in links}

        click_url = cls._safe_hyperlink_address(
            getattr(getattr(shape, "click_action", None), "hyperlink", None)
        )
        if click_url:
            label = PDFExtractor._compact_link_label(
                shape_text or getattr(shape, "name", "") or "Shape link"
            )
            cls._append_link(links, seen, label, click_url, "shape")

        if not getattr(shape, "has_text_frame", False):
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                address = cls._safe_hyperlink_address(getattr(run, "hyperlink", None))
                if not address:
                    continue
                label = PDFExtractor._compact_link_label(run.text or shape_text or "Text link")
                cls._append_link(links, seen, label, address, "run")

    @staticmethod
    def _safe_hyperlink_address(hyperlink: Any) -> str:
        try:
            return str(hyperlink.address or "").strip()
        except Exception:
            return ""

    @staticmethod
    def _append_link(
        links: list[dict[str, str]],
        seen: set[tuple[str, str]],
        label: str,
        target: str,
        kind: str,
    ) -> None:
        key = (label.lower(), target)
        if target and key not in seen:
            seen.add(key)
            links.append({"label": label or "Unlabeled link", "target": target, "kind": kind})

    @staticmethod
    def _table_text(shape: Any) -> str:
        rows = []
        for row in shape.table.rows:
            cells = [PDFExtractor._post_process(cell.text) for cell in row.cells]
            rows.append(" | ".join(cell for cell in cells if cell))
        text = "\n".join(row for row in rows if row.strip())
        return f"Table:\n{text}" if text else ""

    @staticmethod
    def _notes_text(slide: Any) -> str:
        try:
            text_frame = slide.notes_slide.notes_text_frame
        except Exception:
            return ""
        return PDFExtractor._post_process(text_frame.text or "")

    @staticmethod
    def _image_ocr(shape: Any) -> str:
        if not hasattr(shape, "image"):
            return ""
        try:
            from io import BytesIO

            import pytesseract
            from PIL import Image

            image = Image.open(BytesIO(shape.image.blob))
            text = pytesseract.image_to_string(image)
            return PDFExtractor._post_process(text)
        except Exception:
            return ""

    @staticmethod
    def _convert_legacy_ppt(path: Path, output_dir: Path) -> Path | None:
        import shutil
        import subprocess

        executable = shutil.which("soffice") or shutil.which("libreoffice")
        if not executable:
            return None

        try:
            subprocess.run(
                [
                    executable,
                    "--headless",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    str(output_dir),
                    str(path),
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=90,
            )
        except Exception:
            return None

        converted = output_dir / f"{path.stem}.pptx"
        if not converted.exists():
            return None
        return converted


class ExcelExtractor:
    def extract(self, path: Path) -> list[dict[str, Any]]:
        import pandas as pd

        items: list[dict[str, Any]] = []
        sheets = (
            pd.read_excel(path, sheet_name=None, header=None)
            if path.suffix.lower() != ".csv"
            else {"Sheet1": pd.read_csv(path, header=None)}
        )
        for sheet_name, raw_df in sheets.items():
            table_start = self._detect_table_start(raw_df)
            df, header_mode = self._materialize_table(raw_df, table_start)
            df = self._normalize(df)
            profile = self._profile(df)
            preview = df.head(30).to_string(index=False)
            items.append(
                {
                    "text": (
                        f"Sheet: {sheet_name}\n"
                        f"Header detection: {header_mode}\n"
                        f"Schema: {profile}\nPreview:\n{preview}"
                    ),
                    "sheet": sheet_name,
                    "chunk_type": "excel_sheet",
                    "table": df.to_json(orient="records", date_format="iso"),
                }
            )
        return items

    @classmethod
    def _detect_table_start(cls, df) -> int:
        best_index = 0
        best_score = -1.0
        for idx, row in df.head(30).iterrows():
            non_null = row.notna().sum()
            if non_null < 2:
                continue
            following = df.iloc[int(idx) + 1 : int(idx) + 6]
            continuity = 0
            for _, next_row in following.iterrows():
                if next_row.notna().sum() >= max(2, int(non_null * 0.5)):
                    continuity += 1
            score = float(non_null) + continuity * 1.5 + cls._header_likelihood(row, following)
            if score > best_score:
                best_index = int(idx)
                best_score = score
        return best_index

    @classmethod
    def _materialize_table(cls, raw_df, table_start: int):
        table = raw_df.iloc[table_start:].dropna(how="all").dropna(axis=1, how="all")
        if table.empty:
            return table, "empty table"

        candidate_header = table.iloc[0]
        following = table.iloc[1:8]
        if cls._row_looks_like_header(candidate_header, following):
            df = table.iloc[1:].copy()
            df.columns = [
                str(value).strip() if cls._has_value(value) else f"column_{idx + 1}"
                for idx, value in enumerate(candidate_header)
            ]
            return df, f"detected header row {table_start + 1}"

        df = table.copy()
        df.columns = cls._generated_column_names(df)
        return df, f"no reliable header detected; generated {len(df.columns)} column names"

    @classmethod
    def _row_looks_like_header(cls, row, following) -> bool:
        values = [value for value in row.tolist() if cls._has_value(value)]
        if len(values) < 2:
            return False

        kinds = [cls._cell_kind(value) for value in values]
        data_like = sum(kind in {"number", "date", "code", "flag"} for kind in kinds) / len(values)
        if data_like >= 0.35:
            return False

        strings = [str(value).strip() for value in values]
        unique_ratio = len({value.lower() for value in strings}) / max(1, len(strings))
        label_like = sum(cls._looks_like_header_label(value) for value in strings) / len(strings)
        keyword_hits = sum(cls._header_keyword_hit(value) for value in strings)

        comparable = 0
        mismatches = 0
        for column_index, first_value in row.items():
            if not cls._has_value(first_value):
                continue
            next_values = [value for value in following[column_index].tolist() if cls._has_value(value)]
            if not next_values:
                continue
            comparable += 1
            first_kind = cls._cell_kind(first_value)
            next_kinds = [cls._cell_kind(value) for value in next_values[:5]]
            majority_kind = max(set(next_kinds), key=next_kinds.count)
            if first_kind != majority_kind:
                mismatches += 1
        mismatch_ratio = mismatches / max(1, comparable)

        if unique_ratio < 0.8:
            return False
        if keyword_hits >= 2 and label_like >= 0.6:
            return True
        return label_like >= 0.75 and mismatch_ratio >= 0.45

    @classmethod
    def _header_likelihood(cls, row, following) -> float:
        return 2.0 if cls._row_looks_like_header(row, following) else 0.0

    @classmethod
    def _generated_column_names(cls, df) -> list[str]:
        names = []
        counts: dict[str, int] = {}
        for idx, column in enumerate(df.columns):
            series = df[column]
            kind = cls._dominant_column_kind(series)
            base = {
                "date": "date",
                "number": "number",
                "flag": "flag",
                "code": "code",
                "id": "id",
                "sku": "product_sku",
                "status": "status",
            }.get(kind, "text")
            count = counts.get(base, 0) + 1
            counts[base] = count
            names.append(f"{base}_{count}")
        return names

    @classmethod
    def _dominant_column_kind(cls, series) -> str:
        values = [value for value in series.dropna().tolist()[:50] if cls._has_value(value)]
        if not values:
            return "text"
        kinds = [cls._cell_kind(value) for value in values]
        sku_hits = sum(cls._looks_like_sku(value) for value in values)
        status_hits = sum(str(value).strip().upper() in {"OPEN", "CLOSED", "PENDING", "CANCELLED", "ACTIVE", "INACTIVE"} for value in values)
        id_hits = sum(
            bool(re.match(r"^[A-Z]{1,4}[-_/]?\d{2,}$", str(value).strip().upper()))
            for value in values
        )
        if sku_hits / len(values) >= 0.5:
            return "sku"
        if status_hits / len(values) >= 0.5:
            return "status"
        if id_hits / len(values) >= 0.5:
            return "id"
        return max(set(kinds), key=kinds.count)

    @classmethod
    def _cell_kind(cls, value) -> str:
        text = str(value).strip()
        if not text:
            return "blank"
        upper = text.upper()
        if upper in {"Y", "N", "YES", "NO", "TRUE", "FALSE"}:
            return "flag"
        if cls._looks_like_number(value):
            return "number"
        if cls._looks_like_date(value):
            return "date"
        if cls._looks_like_sku(value) or re.match(r"^[A-Z]{2,}[-_/]?[A-Z0-9]{2,}$", upper):
            return "code"
        return "text"

    @staticmethod
    def _has_value(value) -> bool:
        return not (value is None or str(value).strip().lower() in {"", "nan", "none"})

    @staticmethod
    def _looks_like_number(value) -> bool:
        if isinstance(value, (int, float)) and not isinstance(value, bool):
            return True
        text = str(value).strip().replace(",", "").replace("$", "").replace("₹", "").replace("%", "")
        return bool(re.fullmatch(r"[-+]?\d+(?:\.\d+)?", text))

    @staticmethod
    def _looks_like_date(value) -> bool:
        if hasattr(value, "year") and hasattr(value, "month") and hasattr(value, "day"):
            return True
        text = str(value).strip()
        return bool(
            re.fullmatch(r"\d{4}[-/]\d{1,2}[-/]\d{1,2}", text)
            or re.fullmatch(r"\d{1,2}[-/]\d{1,2}[-/]\d{2,4}", text)
        )

    @staticmethod
    def _looks_like_sku(value) -> bool:
        return bool(re.search(r"\bSKU[-_/]?[A-Z0-9]+\b", str(value).strip().upper()))

    @classmethod
    def _looks_like_header_label(cls, value: str) -> bool:
        value = value.strip()
        if not value or len(value) > 80:
            return False
        if cls._looks_like_number(value) or cls._looks_like_date(value) or cls._looks_like_sku(value):
            return False
        if re.search(r"[.!?]{1,}$", value):
            return False
        return bool(re.search(r"[A-Za-z]", value))

    @staticmethod
    def _header_keyword_hit(value: str) -> bool:
        words = set(re.findall(r"[a-z]+", value.lower()))
        keywords = {
            "id",
            "date",
            "name",
            "region",
            "code",
            "sku",
            "product",
            "category",
            "unit",
            "units",
            "qty",
            "quantity",
            "price",
            "cost",
            "discount",
            "status",
            "target",
            "amount",
            "revenue",
            "margin",
            "return",
            "approved",
        }
        return bool(words & keywords)

    @staticmethod
    def _normalize(df):
        df = df.dropna(how="all").dropna(axis=1, how="all")
        columns = [
            re.sub(r"_+", "_", re.sub(r"[^a-z0-9]+", "_", str(col).strip().lower())).strip("_")
            or f"column_{idx}"
            for idx, col in enumerate(df.columns)
        ]
        seen: dict[str, int] = {}
        unique_columns = []
        for column in columns:
            seen[column] = seen.get(column, 0) + 1
            unique_columns.append(column if seen[column] == 1 else f"{column}_{seen[column]}")
        df.columns = unique_columns
        return df

    @staticmethod
    def _profile(df) -> dict[str, str]:
        return {column: str(dtype) for column, dtype in df.dtypes.items()}


class ChartOCRExtractor:
    def extract(self, path: Path) -> list[dict[str, Any]]:
        import cv2
        import pytesseract

        image = cv2.imread(str(path))
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
        text = pytesseract.image_to_string(binary)
        return [
            {
                "text": "Chart/image OCR observations:\n" + PDFExtractor._post_process(text),
                "chunk_type": "chart_ocr",
                "confidence": 0.55,
            }
        ]
