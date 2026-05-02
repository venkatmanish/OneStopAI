from pathlib import Path
from types import SimpleNamespace
import json

from core.ingestion import IngestionService, PDFExtractor, PowerPointExtractor, chunk_text


def test_pdf_extractor_indexes_hidden_hyperlink_target(tmp_path: Path):
    import fitz

    pdf_path = tmp_path / "resume.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "GitHub LinkedIn")
    github_rect = page.search_for("GitHub")[0]
    page.insert_link(
        {
            "kind": fitz.LINK_URI,
            "from": github_rect,
            "uri": "https://github.com/manish-ai",
        }
    )
    doc.save(pdf_path)
    doc.close()

    items = PDFExtractor().extract(pdf_path)
    link_text = "\n".join(item["text"] for item in items if item.get("chunk_type") == "pdf_links")

    assert "GitHub -> https://github.com/manish-ai" in link_text


def test_powerpoint_extractor_indexes_slide_text_and_hyperlinks(tmp_path: Path):
    from pptx import Presentation

    pptx_path = tmp_path / "portfolio.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(914400, 914400, 4572000, 914400)
    paragraph = textbox.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "GitHub"
    run.hyperlink.address = "https://github.com/manish-ai"
    presentation.save(pptx_path)

    items = PowerPointExtractor().extract(pptx_path)
    text = "\n".join(item["text"] for item in items)

    assert "Slide 1" in text
    assert "GitHub" in text
    assert "GitHub -> https://github.com/manish-ai" in text


def test_ingestion_reindexes_chunks_from_older_parser_version():
    old_chunk = SimpleNamespace(extra={"parser_version": "old"})
    current_chunk = SimpleNamespace(
        extra={
            "parser_version": IngestionService.PARSER_VERSION,
            "embedding_model": IngestionService().embedding_model.settings.embedding_model,
        }
    )
    old_embedding_chunk = SimpleNamespace(
        extra={
            "parser_version": IngestionService.PARSER_VERSION,
            "embedding_model": "old-embedding-model",
        }
    )

    assert IngestionService._needs_reindex([old_chunk]) is True
    assert IngestionService._needs_reindex([current_chunk]) is False
    assert IngestionService._needs_reindex([old_embedding_chunk]) is True


def test_chunk_text_prefers_semantic_boundaries():
    text = (
        "Section A\n"
        "This paragraph explains the first idea in detail. It should stay together.\n\n"
        "Section B\n"
        "This paragraph explains the second idea in detail. It should also stay together.\n\n"
        "Section C\n"
        "This paragraph explains the third idea in detail."
    )

    chunks = chunk_text(text, size=120, overlap=35)

    assert len(chunks) > 1
    assert all(chunk.strip() for chunk in chunks)
    assert any("Section B" in chunk for chunk in chunks)


def test_excel_extractor_does_not_require_optional_tabulate(tmp_path: Path):
    from openpyxl import Workbook

    from core.ingestion import ExcelExtractor

    workbook_path = tmp_path / "messy.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Messy"
    ws["A1"] = "noise row before headers"
    ws.append([])
    ws.append(["Region", "Revenue", "Cost"])
    ws.append(["South", 1000, 600])
    ws.append(["North", 800, 500])
    wb.save(workbook_path)

    items = ExcelExtractor().extract(workbook_path)

    assert items
    assert "Schema:" in items[0]["text"]
    assert "South" in items[0]["text"]


def test_excel_extractor_preserves_headerless_first_data_row(tmp_path: Path):
    from openpyxl import Workbook

    from core.ingestion import ExcelExtractor

    workbook_path = tmp_path / "headerless.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "NoHeader"
    ws.append(["O-1001", "SOUTH", "SKU-ALPHA", 12, 1450, "Closed"])
    ws.append(["O-1002", "WEST", "SKU-BETA", 7, 2300, "Closed"])
    wb.save(workbook_path)

    item = ExcelExtractor().extract(workbook_path)[0]
    records = json.loads(item["table"])

    assert "no reliable header detected" in item["text"]
    assert records[0]["id_1"] == "O-1001"
    assert records[0]["code_1"] == "SOUTH"
    assert records[0]["product_sku_1"] == "SKU-ALPHA"
    assert records[0]["number_1"] == 12
    assert records[0]["status_1"] == "Closed"


def test_chunk_rows_store_parent_context():
    service = IngestionService()
    rows = service._build_chunk_rows(
        "doc",
        "version",
        "sample.pdf",
        [
            {
                "text": "Heading\n\n" + "A detailed paragraph about retrieval quality. " * 80,
                "page": 1,
                "chunk_type": "pdf_text",
            }
        ],
        "upload",
    )

    assert len(rows) > 1
    assert rows[0]["extra"]["parent_context"].startswith("Heading")
    assert rows[0]["extra"]["parent_id"] == rows[1]["extra"]["parent_id"]
    assert rows[0]["extra"]["chunk_count"] == len(rows)
    assert rows[0]["extra"]["embedding_model"] == service.embedding_model.settings.embedding_model
